"""PowerPoint presentation generator."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from spa.models import PresentationContent, SlideContent


class PPTXGenerator:
    """Generator for PowerPoint presentations."""

    # Brand colors
    PRIMARY_COLOR = RgbColor(0x00, 0x52, 0xCC)  # Blue
    SECONDARY_COLOR = RgbColor(0x33, 0x33, 0x33)  # Dark gray
    ACCENT_COLOR = RgbColor(0x00, 0x96, 0x88)  # Teal

    def __init__(self) -> None:
        """Initialize the generator."""
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def generate(self, content: PresentationContent, output_path: Path) -> Path:
        """Generate PowerPoint presentation from content."""
        for slide_content in content.slides:
            self._add_slide(slide_content)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        return output_path

    def _add_slide(self, content: SlideContent) -> None:
        """Add a slide based on content type."""
        if content.slide_type == "title":
            self._add_title_slide(content)
        else:
            self._add_content_slide(content)

    def _add_title_slide(self, content: SlideContent) -> None:
        """Add a title slide."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = content.title
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = self.PRIMARY_COLOR
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle / Main message
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(12.333), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = content.main_message
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = self.SECONDARY_COLOR
        subtitle_para.alignment = PP_ALIGN.CENTER

        self._add_speaker_notes(slide, content.speaker_notes)

    def _add_content_slide(self, content: SlideContent) -> None:
        """Add a content slide with title and bullet points."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Title bar background
        title_bg = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(13.333), Inches(1.2)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = self.PRIMARY_COLOR
        title_bg.line.fill.background()

        # Title text
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = content.title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RgbColor(0xFF, 0xFF, 0xFF)

        # Main message
        message_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(12.333), Inches(1)
        )
        message_frame = message_box.text_frame
        message_para = message_frame.paragraphs[0]
        message_para.text = content.main_message
        message_para.font.size = Pt(20)
        message_para.font.bold = True
        message_para.font.color.rgb = self.SECONDARY_COLOR

        # Bullet points
        if content.bullet_points:
            bullets_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(2.5), Inches(11.933), Inches(4)
            )
            bullets_frame = bullets_box.text_frame
            bullets_frame.word_wrap = True

            for i, point in enumerate(content.bullet_points):
                if i == 0:
                    para = bullets_frame.paragraphs[0]
                else:
                    para = bullets_frame.add_paragraph()
                para.text = f"• {point}"
                para.font.size = Pt(18)
                para.font.color.rgb = self.SECONDARY_COLOR
                para.space_before = Pt(12)

        self._add_speaker_notes(slide, content.speaker_notes)

    def _add_speaker_notes(self, slide, notes: str) -> None:
        """Add speaker notes to a slide."""
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes
